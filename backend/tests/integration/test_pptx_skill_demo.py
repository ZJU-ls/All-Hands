"""Integration · build a real 12-page deck out of the skill's templates.

Pins the end-to-end shape of the design skill: templates are loaded
verbatim from the skill's `templates/` directory, content fields are
overridden per slide, and the renderer produces a valid Office-
compatible .pptx the user can open in PowerPoint / Keynote / WPS.

This test exists to catch:
- a template stops parsing (someone hand-edited JSON badly)
- a template references a primitive shape the renderer no longer
  accepts (schema drift)
- the deck produced by the templates is silently empty / corrupt
- the chart template produces a chart shape after round-trip

If you change a template's outer shape (fields the LLM overrides),
make sure the override block in this test still aligns.
"""

from __future__ import annotations

import io
import json
from copy import deepcopy
from pathlib import Path
from typing import Any

import pytest

from allhands.execution.artifact_generators.pptx import (
    extract_slide_text,
    render_pptx,
)

TEMPLATES = (
    Path(__file__).resolve().parents[2] / "skills" / "builtin" / "allhands-design" / "templates"
)


def _load(name: str) -> dict[str, Any]:
    spec = json.loads((TEMPLATES / name).read_text(encoding="utf-8"))
    spec.pop("_doc", None)
    return spec


def _set_text(spec: dict[str, Any], shape_index: int, text: str) -> None:
    """Replace the `text` field on a specific shape (LLM does this in chat)."""
    spec["shapes"][shape_index]["text"] = text


def _build_demo_deck() -> list[dict[str, Any]]:
    """Twelve slides drawn from the dark-theme template set:

    1.  cover
    2.  agenda
    3.  section-divider (chapter 1)
    4.  title-content
    5.  title-content
    6.  section-divider (chapter 2)
    7.  two-column
    8.  image-text
    9.  kpi-grid
    10. chart-with-caption
    11. quote
    12. closing
    """
    cover = _load("cover-dark.json")
    _set_text(cover, 2, "ALL HANDS · v0")
    _set_text(cover, 3, "All Hands\nfor your team")
    _set_text(cover, 4, "Design / dispatch / observe a team of AI employees")
    _set_text(cover, 5, "AUTHOR · 2026-04-27")

    agenda = _load("agenda-dark.json")

    s1 = _load("section-divider-dark.json")
    _set_text(s1, 1, "01")
    _set_text(s1, 3, "Why we built this")
    _set_text(s1, 4, "Replacing 30-tool sprawl with one chat")

    p1 = _load("title-content-dark.json")
    _set_text(p1, 2, "The problem")
    _set_text(p1, 4, "Most teams string together 30+ point tools.")
    _set_text(
        p1,
        5,
        "All Hands collapses that surface into one platform that any team "
        "member can drive from chat — without learning a new admin UI.",
    )

    p2 = _load("title-content-dark.json")
    _set_text(p2, 2, "Our wedge")
    _set_text(p2, 4, "Tool First · everything is a tool · chat-driven config")
    _set_text(
        p2,
        5,
        "The same operations available in the UI are exposed as Meta Tools, "
        "so the Lead Agent can do every admin action via conversation.\n\n"
        "Skills inject capabilities on demand · the model only pays the "
        "context cost when the work calls for it.",
    )

    s2 = _load("section-divider-dark.json")
    _set_text(s2, 1, "02")
    _set_text(s2, 3, "How it shows up")
    _set_text(s2, 4, "Three concrete patterns end users will feel")

    two = _load("two-column-dark.json")
    _set_text(two, 2, "Before & after")
    _set_text(two, 4, "Before · 30 tools")
    _set_text(
        two,
        5,
        "Each tool needs its own login · API key · admin doc.\n\n"
        "On-call rotates between dashboards.\n\n"
        "Most actions are out of reach for non-engineers.",
    )
    _set_text(two, 7, "After · one chat")
    _set_text(
        two,
        8,
        "All operations are tools the Lead Agent can drive.\n\n"
        "On-call asks one assistant.\n\n"
        "Anyone on the team can do real work end-to-end.",
    )

    img = _load("image-text-dark.json")
    _set_text(img, 2, "Lead Agent in action")
    _set_text(img, 5, "PRODUCT SCREENSHOT")
    _set_text(
        img,
        6,
        "•  Single chat surface\n\n"
        "•  Plan + tool-call streaming\n\n"
        "•  Native artifact panel\n\n"
        "•  Skill descriptors always visible",
    )

    kpi = _load("kpi-grid-dark.json")
    _set_text(kpi, 2, "Numbers worth quoting")
    _set_text(kpi, 3, "From design partner alpha")

    chart = _load("chart-with-caption-dark.json")
    _set_text(chart, 2, "Adoption curve")
    _set_text(chart, 3, "Active conversations grew 3x quarter over quarter")
    chart_shape = chart["shapes"][4]
    chart_shape["categories"] = ["Q1", "Q2", "Q3", "Q4"]
    chart_shape["series"] = [
        {"name": "All Hands", "values": [40, 95, 188, 360], "color_hex": "#3b82f6"},
        {"name": "Industry avg", "values": [45, 88, 130, 180], "color_hex": "#60a5fa"},
    ]

    q = _load("quote-dark.json")
    _set_text(
        q,
        2,
        "“The day a team replaces six dashboards with one chat is the day they really use AI.”",
    )
    _set_text(q, 4, "— Internal design review · 2026-04")

    close = _load("closing-dark.json")
    # closing-dark.json shape indices:
    # 0 left_strip · 1 hairline · 2 eyebrow · 3 headline · 4 accent rect ·
    # 5 NEXT STEPS eyebrow · 6 body bullets · 7 footer
    _set_text(close, 3, "Want to try All Hands with your team?")
    _set_text(
        close,
        6,
        "1.  Pilot enrollment by 2026-05-15\n\n"
        "2.  20-min walkthrough · book at /demo\n\n"
        "3.  Decision needed: pick design partner cohort by EOM",
    )
    _set_text(close, 7, "PRODUCT · pilot@allhands.io")

    return [
        cover,
        agenda,
        s1,
        p1,
        p2,
        s2,
        two,
        img,
        kpi,
        chart,
        q,
        close,
    ]


def test_skill_templates_compose_into_a_real_deck(tmp_path: Path) -> None:
    slides = _build_demo_deck()
    blob, warnings = render_pptx(
        page={"background": {"color_hex": "#0a0e1a"}},
        slides=slides,
    )
    assert warnings == []
    # > 30 KB · plan §6.2 acceptance threshold
    assert len(blob) > 30_000

    # extract_slide_text round trip · 12 slides
    outline = extract_slide_text(blob)
    assert len(outline) == 12

    # cover title made it through · plan §6.2 acceptance ("All Hands" appears
    # in slide 1 title)
    assert "ALL HANDS" in outline[0]["title"] or "All Hands" in outline[0]["title"]

    # chart slide can be re-opened and a chart shape is present
    from pptx import Presentation

    parsed = Presentation(io.BytesIO(blob))
    assert len(parsed.slides) == 12
    chart_slide = parsed.slides[9]  # 0-indexed · slot 10 in the spec
    assert any(s.has_chart for s in chart_slide.shapes), "chart shape missing"

    # Optional sanity: deep-copy isolation (mutating the spec we built
    # shouldn't reach back into the on-disk template)
    original_cover = json.loads((TEMPLATES / "cover-dark.json").read_text())
    assert original_cover["shapes"][2]["text"] == "EYEBROW"


def test_light_theme_templates_also_render() -> None:
    """Quick parity smoke · the light variants should render too."""
    light_specs = []
    for name in ("cover-light.json", "title-content-light.json", "closing-light.json"):
        spec = json.loads((TEMPLATES / name).read_text(encoding="utf-8"))
        spec.pop("_doc", None)
        light_specs.append(deepcopy(spec))
    blob, warnings = render_pptx(
        page={"background": {"color_hex": "#ffffff"}},
        slides=light_specs,
    )
    assert warnings == []
    assert len(blob) > 5_000


def test_templates_each_render_individually() -> None:
    """Each of the 20 templates must render as a single-slide deck."""
    paths = sorted(TEMPLATES.glob("*.json"))
    assert len(paths) == 20, f"expected 20 templates, got {len(paths)}"
    for p in paths:
        spec = json.loads(p.read_text(encoding="utf-8"))
        spec.pop("_doc", None)
        blob, warnings = render_pptx(slides=[spec])
        assert warnings == [], f"{p.name}: warnings={warnings}"
        assert len(blob) > 1_000, f"{p.name}: blob too small ({len(blob)})"


@pytest.mark.parametrize("theme", ["dark", "light"])
def test_template_pair_has_same_shape_count(theme: str) -> None:
    """The dark and light variants must mirror each other so swapping
    themes doesn't drop shapes."""
    base_names = (
        "cover",
        "agenda",
        "section-divider",
        "title-content",
        "two-column",
        "image-text",
        "kpi-grid",
        "chart-with-caption",
        "quote",
        "closing",
    )
    other_theme = "light" if theme == "dark" else "dark"
    for name in base_names:
        a = json.loads((TEMPLATES / f"{name}-{theme}.json").read_text())
        b = json.loads((TEMPLATES / f"{name}-{other_theme}.json").read_text())
        assert len(a["shapes"]) == len(b["shapes"]), (
            f"{name}: {theme}={len(a['shapes'])} shapes vs {other_theme}={len(b['shapes'])}"
        )
