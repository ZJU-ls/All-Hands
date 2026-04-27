"""Unit tests for the primitives-mode pptx renderer.

Covers the `text` / `rect` / `line` / `image` / `chart` builders end to
end (single-shape slide → bytes → re-parsed via python-pptx) and
boundary cases that must surface as ToolArgError so the pipeline
emits ADR 0021 envelopes.

The renderer's contract: caller supplies all visual choices, the
output is a valid Office-compatible .pptx; errors are educational.
"""

from __future__ import annotations

import base64
import io
from pathlib import Path
from typing import Any

import pytest

from allhands.execution.artifact_generators.pdf import ArtifactGenerationError
from allhands.execution.artifact_generators.pptx import (
    extract_slide_text,
    render_pptx,
)
from allhands.execution.tool_arg_validation import ToolArgError


def _open(blob: bytes) -> Any:
    from pptx import Presentation

    return Presentation(io.BytesIO(blob))


# A 1x1 transparent PNG · 67 bytes · base64.
_TINY_PNG_B64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
    "+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
)


# ---- happy paths · single primitive per slide ------------------------------


def test_text_primitive_round_trips() -> None:
    blob, warnings = render_pptx(
        slides=[
            {
                "shapes": [
                    {
                        "type": "text",
                        "x": 1,
                        "y": 1,
                        "w": 6,
                        "h": 1,
                        "text": "Hello, primitives.",
                        "font": {
                            "family": "Inter",
                            "size": 32,
                            "weight": "bold",
                            "color_hex": "#0a0e1a",
                        },
                        "align": "center",
                        "v_align": "middle",
                    }
                ]
            }
        ],
    )
    assert isinstance(blob, bytes) and len(blob) > 1024
    assert warnings == []
    parsed = _open(blob)
    assert len(parsed.slides) == 1
    extracted = extract_slide_text(blob)
    assert extracted == [{"title": "Hello, primitives.", "body": []}]


def test_rect_with_corner_radius_and_border() -> None:
    blob, _ = render_pptx(
        slides=[
            {
                "shapes": [
                    {
                        "type": "rect",
                        "x": 0.5,
                        "y": 0.5,
                        "w": 12,
                        "h": 1,
                        "fill_hex": "#3b82f6",
                        "border": {"color_hex": "#1e293b", "thickness_pt": 1.5},
                        "corner_radius_in": 0.2,
                    }
                ]
            }
        ],
    )
    parsed = _open(blob)
    slide = parsed.slides[0]
    # one rect shape, no text frame (background full-canvas rect skipped
    # because no background was set)
    assert any(getattr(s, "shape_type", None) is not None for s in slide.shapes)


def test_line_with_dash() -> None:
    blob, _ = render_pptx(
        slides=[
            {
                "shapes": [
                    {
                        "type": "line",
                        "x1": 0.5,
                        "y1": 1.5,
                        "x2": 12.8,
                        "y2": 1.5,
                        "color_hex": "#3b82f6",
                        "thickness_pt": 2,
                        "dash": "dash",
                    }
                ]
            }
        ],
    )
    parsed = _open(blob)
    assert len(parsed.slides) == 1


def test_image_primitive_embeds_png() -> None:
    blob, _ = render_pptx(
        slides=[
            {
                "shapes": [
                    {
                        "type": "image",
                        "x": 5,
                        "y": 2,
                        "w": 3,
                        "h": 3,
                        "data_b64": _TINY_PNG_B64,
                        "alt_text": "demo",
                    }
                ]
            }
        ],
    )
    parsed = _open(blob)
    pictures = [s for s in parsed.slides[0].shapes if s.shape_type == 13]
    assert len(pictures) == 1


def test_image_strips_data_url_prefix() -> None:
    blob, _ = render_pptx(
        slides=[
            {
                "shapes": [
                    {
                        "type": "image",
                        "x": 1,
                        "y": 1,
                        "w": 1,
                        "h": 1,
                        "data_b64": f"data:image/png;base64,{_TINY_PNG_B64}",
                    }
                ]
            }
        ],
    )
    assert len(blob) > 1024


@pytest.mark.parametrize("chart_type", ["bar", "column", "line", "pie"])
def test_chart_primitive_renders_each_type(chart_type: str) -> None:
    blob, _ = render_pptx(
        slides=[
            {
                "shapes": [
                    {
                        "type": "chart",
                        "x": 1,
                        "y": 1,
                        "w": 10,
                        "h": 5,
                        "chart_type": chart_type,
                        "categories": ["Q1", "Q2", "Q3", "Q4"],
                        "series": [
                            {
                                "name": "Revenue",
                                "values": [12, 15, 18, 22],
                                "color_hex": "#3b82f6",
                            }
                        ],
                        "show_legend": True,
                        "show_data_labels": False,
                    }
                ]
            }
        ],
    )
    parsed = _open(blob)
    charts = [s for s in parsed.slides[0].shapes if s.has_chart]
    assert len(charts) == 1


def test_speaker_notes_round_trip() -> None:
    blob, _ = render_pptx(
        slides=[
            {
                "shapes": [
                    {
                        "type": "text",
                        "x": 1,
                        "y": 1,
                        "w": 6,
                        "h": 1,
                        "text": "Cover",
                    }
                ],
                "notes": "Open with the headline · 30 seconds.",
            }
        ],
    )
    parsed = _open(blob)
    nf = parsed.slides[0].notes_slide.notes_text_frame
    assert "Open with the headline" in nf.text


def test_page_size_and_background_apply() -> None:
    blob, _ = render_pptx(
        page={
            "width_in": 16,
            "height_in": 9,
            "background": {"color_hex": "#0a0e1a"},
        },
        slides=[{"shapes": [{"type": "text", "x": 1, "y": 1, "w": 4, "h": 1, "text": "x"}]}],
    )
    parsed = _open(blob)
    # python-pptx stores width / height as EMU. 16 inches = 14_630_400 EMU
    assert parsed.slide_width == 14630400
    assert parsed.slide_height == 8229600
    # Background rect added (full-canvas), so slide has at least 2 shapes.
    assert len(parsed.slides[0].shapes) >= 2


def test_per_slide_background_overrides_page_background() -> None:
    blob, _ = render_pptx(
        page={"background": {"color_hex": "#ffffff"}},
        slides=[
            {
                "background": {"color_hex": "#0a0e1a"},
                "shapes": [{"type": "text", "x": 1, "y": 1, "w": 4, "h": 1, "text": "a"}],
            }
        ],
    )
    parsed = _open(blob)
    # at least the bg rect + the textbox
    assert len(parsed.slides[0].shapes) >= 2


# ---- structured failure path · ADR 0021 envelopes -------------------------


def test_empty_slides_raises_artifact_error() -> None:
    """The empty-deck path is library-level (no caller field), so it
    surfaces as ArtifactGenerationError rather than ToolArgError."""
    with pytest.raises(ArtifactGenerationError):
        render_pptx(slides=[])


def test_slide_without_shapes_raises_tool_arg_error() -> None:
    with pytest.raises(ToolArgError) as ex:
        render_pptx(slides=[{"shapes": []}])
    assert ex.value.field == "slides[0].shapes"
    assert ex.value.expected == "non-empty array"


def test_unknown_shape_type_raises_tool_arg_error() -> None:
    with pytest.raises(ToolArgError) as ex:
        render_pptx(slides=[{"shapes": [{"type": "video"}]}])
    assert ex.value.field == "slides[0].shapes[0].type"
    assert "['text', 'rect', 'line', 'image', 'chart']" in ex.value.expected


def test_bad_color_hex_raises() -> None:
    with pytest.raises(ToolArgError) as ex:
        render_pptx(
            slides=[
                {
                    "shapes": [
                        {
                            "type": "rect",
                            "x": 1,
                            "y": 1,
                            "w": 1,
                            "h": 1,
                            "fill_hex": "blue",
                        }
                    ]
                }
            ]
        )
    assert ex.value.field == "slides[0].shapes[0].fill_hex"
    assert "#RRGGBB" in ex.value.expected


def test_shape_off_canvas_raises_with_explicit_field() -> None:
    with pytest.raises(ToolArgError) as ex:
        render_pptx(
            slides=[
                {
                    "shapes": [
                        {
                            "type": "rect",
                            "x": 12,
                            "y": 1,
                            "w": 5,  # 12 + 5 = 17 > page width 13.333
                            "h": 1,
                            "fill_hex": "#000000",
                        }
                    ]
                }
            ]
        )
    assert ex.value.field == "slides[0].shapes[0].x+w"
    assert "page width" in ex.value.expected


def test_chart_values_length_mismatch_raises_with_hint() -> None:
    with pytest.raises(ToolArgError) as ex:
        render_pptx(
            slides=[
                {
                    "shapes": [
                        {
                            "type": "chart",
                            "x": 1,
                            "y": 1,
                            "w": 8,
                            "h": 4,
                            "chart_type": "bar",
                            "categories": ["a", "b", "c"],
                            "series": [{"name": "s", "values": [1, 2]}],
                        }
                    ]
                }
            ]
        )
    assert ex.value.field == "slides[0].shapes[0].series[0].values"
    assert "length 3" in ex.value.expected
    assert "categories has 3 entries" in ex.value.hint


def test_image_invalid_base64_raises_with_decode_error() -> None:
    with pytest.raises(ToolArgError) as ex:
        render_pptx(
            slides=[
                {
                    "shapes": [
                        {
                            "type": "image",
                            "x": 1,
                            "y": 1,
                            "w": 1,
                            "h": 1,
                            "data_b64": "this is !!! not base64",
                        }
                    ]
                }
            ]
        )
    assert ex.value.field == "slides[0].shapes[0].data_b64"
    assert "valid base64" in ex.value.expected
    assert "PNG" in ex.value.hint or "JPEG" in ex.value.hint


def test_image_oversize_raises() -> None:
    big = base64.b64encode(b"\x00" * (3 * 1024 * 1024)).decode()
    with pytest.raises(ToolArgError) as ex:
        render_pptx(
            slides=[
                {
                    "shapes": [
                        {
                            "type": "image",
                            "x": 1,
                            "y": 1,
                            "w": 1,
                            "h": 1,
                            "data_b64": big,
                        }
                    ]
                }
            ]
        )
    assert "2048KB" in ex.value.expected or "<= 2048" in ex.value.expected
    assert "Compress" in ex.value.hint or "compress" in ex.value.hint.lower()


def test_chart_categories_missing_raises() -> None:
    with pytest.raises(ToolArgError) as ex:
        render_pptx(
            slides=[
                {
                    "shapes": [
                        {
                            "type": "chart",
                            "x": 1,
                            "y": 1,
                            "w": 8,
                            "h": 4,
                            "chart_type": "bar",
                            "series": [{"name": "s", "values": [1]}],
                        }
                    ]
                }
            ]
        )
    assert ex.value.field == "slides[0].shapes[0].categories"
    assert "non-empty array" in ex.value.expected


def test_text_invalid_align_enum_raises() -> None:
    with pytest.raises(ToolArgError) as ex:
        render_pptx(
            slides=[
                {
                    "shapes": [
                        {
                            "type": "text",
                            "x": 1,
                            "y": 1,
                            "w": 1,
                            "h": 1,
                            "text": "x",
                            "align": "justify",
                        }
                    ]
                }
            ]
        )
    assert ex.value.field == "slides[0].shapes[0].align"
    assert "left" in ex.value.expected and "center" in ex.value.expected


# ---- extract_slide_text ----------------------------------------------------


def test_extract_slide_text_picks_largest_text_as_title() -> None:
    blob, _ = render_pptx(
        slides=[
            {
                "shapes": [
                    {
                        "type": "text",
                        "x": 1,
                        "y": 1,
                        "w": 11,
                        "h": 1.2,
                        "text": "Big Headline",
                        "font": {"size": 48},
                    },
                    {
                        "type": "text",
                        "x": 1,
                        "y": 3,
                        "w": 11,
                        "h": 0.8,
                        "text": "small body line one",
                        "font": {"size": 14},
                    },
                ]
            }
        ],
    )
    extracted = extract_slide_text(blob)
    assert extracted[0]["title"] == "Big Headline"
    assert "small body line one" in extracted[0]["body"]


def test_extract_slide_text_returns_empty_on_corrupt_blob() -> None:
    assert extract_slide_text(b"not a real pptx") == []


# ---- regression bytes get written to disk reasonably -----------------------


def test_renderer_writes_compact_blob_for_simple_deck(tmp_path: Path) -> None:
    blob, _ = render_pptx(
        slides=[
            {"shapes": [{"type": "text", "x": 1, "y": 1, "w": 4, "h": 1, "text": "1"}]},
            {"shapes": [{"type": "text", "x": 1, "y": 1, "w": 4, "h": 1, "text": "2"}]},
        ],
    )
    out = tmp_path / "demo.pptx"
    out.write_bytes(blob)
    assert out.stat().st_size > 1024
    # Re-open via python-pptx as a sanity round-trip.
    re_parsed = _open(out.read_bytes())
    assert len(re_parsed.slides) == 2
