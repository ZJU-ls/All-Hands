"""PPTX generator · primitive shape spec → .pptx bytes via python-pptx.

The renderer is **opinion-free**: every visual choice (color, position,
font, sizes) comes from the caller. Layout templates and design language
live in skill files outside this module — see ADR 0021 for the contract
boundary.

Five primitive shapes are supported (no further extensions; expanding
the surface is a tool-contract change, not a renderer change):

    text   absolute-positioned text frame · color / weight / align
    rect   axis-aligned rectangle · fill + border + corner radius
    line   straight line · color / thickness / dash
    image  base64 PNG/JPEG embed · no network fetch
    chart  native python-pptx chart · bar / line / pie / column

Validation errors raise `ToolArgError` so the tool pipeline wraps them
into ADR 0021 structured envelopes (field/expected/received/hint).
Library / sandbox failures raise `ArtifactGenerationError` and reach
the executor as plain `{error: str(exc)}`.
"""

from __future__ import annotations

import base64
import binascii
import io
import re
from typing import Any

from allhands.execution.artifact_generators.pdf import ArtifactGenerationError
from allhands.execution.tool_arg_validation import ToolArgError

_HEX_RE = re.compile(r"^#[0-9a-fA-F]{6}$")
_DEFAULT_PAGE_W = 13.333  # 16:9 widescreen
_DEFAULT_PAGE_H = 7.5
_MIN_DIM_IN = 1.0
_MAX_DIM_IN = 56.0
_MAX_IMAGE_BYTES = 2 * 1024 * 1024
_DEFAULT_FONT_PT = 18
_VALID_PRIMITIVES = ("text", "rect", "line", "image", "chart")
_VALID_CHART_TYPES = ("bar", "line", "pie", "column")
_VALID_DASH = ("solid", "dash", "dot")
_VALID_ALIGN = ("left", "center", "right")
_VALID_V_ALIGN = ("top", "middle", "bottom")
_VALID_WEIGHT = ("normal", "bold")


def _bad(field: str, expected: str, received: str, hint: str) -> ToolArgError:
    return ToolArgError(field=field, expected=expected, received=received, hint=hint)


def _summary(value: Any) -> str:
    if value is None:
        return "null"
    if isinstance(value, bool):
        return f"boolean {value}"
    if isinstance(value, int | float):
        return f"number {value}"
    if isinstance(value, str):
        v = value if len(value) <= 60 else value[:57] + "..."
        return f"string {v!r}"
    if isinstance(value, list):
        return f"array (length {len(value)})"
    if isinstance(value, dict):
        return f"object with keys {list(value.keys())[:5]}"
    return type(value).__name__


def _as_dict(value: Any, field: str) -> dict[str, Any]:
    if not isinstance(value, dict):
        raise _bad(
            field=field,
            expected="object",
            received=_summary(value),
            hint=f"`{field}` must be a JSON object literal.",
        )
    return value


def _check_hex(field: str, value: Any) -> str:
    if not isinstance(value, str) or not _HEX_RE.match(value):
        raise _bad(
            field=field,
            expected="hex color string '#RRGGBB'",
            received=_summary(value),
            hint=f"Set `{field}` to e.g. '#3b82f6' (six hex digits with leading hash).",
        )
    return value


def _check_number(field: str, value: Any, *, lo: float, hi: float) -> float:
    if isinstance(value, bool) or not isinstance(value, int | float):
        raise _bad(
            field=field,
            expected="number",
            received=_summary(value),
            hint=f"`{field}` must be a JSON number (no quotes).",
        )
    if value < lo or value > hi:
        raise _bad(
            field=field,
            expected=f"number in [{lo}, {hi}]",
            received=_summary(value),
            hint=f"`{field}` is in inches; keep it within page bounds.",
        )
    return float(value)


def _check_enum(field: str, value: Any, choices: tuple[str, ...]) -> str:
    if value not in choices:
        raise _bad(
            field=field,
            expected=f"one of {list(choices)}",
            received=_summary(value),
            hint=f"Set `{field}` to one of {list(choices)}.",
        )
    return str(value)


# ---- python-pptx primitives -----------------------------------------------


def _rgb(hex_str: str) -> Any:
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(hex_str.lstrip("#"))  # type: ignore[no-untyped-call]


def _emu_inch(inches: float) -> Any:
    from pptx.util import Inches

    return Inches(inches)


def _set_slide_background(slide: Any, color_hex: str) -> None:
    """Force a solid-fill background spanning the whole slide.

    python-pptx's `slide.background` API has limited cross-version
    reliability for fill changes; a full-canvas rectangle is the most
    portable approach. The rectangle is `send_to_back` so it sits under
    every other shape.
    """
    from pptx.enum.shapes import MSO_SHAPE

    width = slide.part.package.presentation_part.presentation.slide_width
    height = slide.part.package.presentation_part.presentation.slide_height
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(color_hex)
    bg.line.fill.background()
    # Move to the bottom of the z-order.
    sp_tree = bg._element.getparent()
    sp_tree.remove(bg._element)
    sp_tree.insert(2, bg._element)


def _shape_bounds(
    field: str, shape: dict[str, Any], page_w: float, page_h: float
) -> tuple[float, float, float, float]:
    x = _check_number(f"{field}.x", shape.get("x"), lo=0, hi=page_w)
    y = _check_number(f"{field}.y", shape.get("y"), lo=0, hi=page_h)
    w = _check_number(f"{field}.w", shape.get("w"), lo=0.01, hi=page_w)
    h = _check_number(f"{field}.h", shape.get("h"), lo=0.01, hi=page_h)
    if x + w > page_w + 1e-6:
        raise _bad(
            field=f"{field}.x+w",
            expected=f'x + w <= page width ({page_w}")',
            received=f"x={x} w={w} sum={x + w}",
            hint="Reduce x or w so the shape fits within the page.",
        )
    if y + h > page_h + 1e-6:
        raise _bad(
            field=f"{field}.y+h",
            expected=f'y + h <= page height ({page_h}")',
            received=f"y={y} h={h} sum={y + h}",
            hint="Reduce y or h so the shape fits within the page.",
        )
    return x, y, w, h


def _render_text(
    slide: Any, field: str, shape: dict[str, Any], page_w: float, page_h: float
) -> None:
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Pt

    x, y, w, h = _shape_bounds(field, shape, page_w, page_h)
    text_value = shape.get("text")
    if not isinstance(text_value, str):
        raise _bad(
            field=f"{field}.text",
            expected="string",
            received=_summary(text_value),
            hint="Pass `text` as a string (use \\n for line breaks).",
        )

    box = slide.shapes.add_textbox(_emu_inch(x), _emu_inch(y), _emu_inch(w), _emu_inch(h))
    tf = box.text_frame
    tf.word_wrap = True
    # Default no padding so layouts line up to the inch grid.
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    font_spec = shape.get("font")
    family: str | None = None
    size_pt: int = _DEFAULT_FONT_PT
    weight = "normal"
    italic = False
    color: str | None = None
    if font_spec is not None:
        font_d = _as_dict(font_spec, f"{field}.font")
        if "family" in font_d:
            fam = font_d["family"]
            if not isinstance(fam, str) or not fam.strip():
                raise _bad(
                    field=f"{field}.font.family",
                    expected="non-empty string",
                    received=_summary(fam),
                    hint="Pass a font family name (e.g. 'Helvetica') · the system falls back if not installed.",
                )
            family = fam
        if "size" in font_d:
            size_pt_raw = font_d["size"]
            if (
                not isinstance(size_pt_raw, int | float)
                or isinstance(size_pt_raw, bool)
                or size_pt_raw < 1
            ):
                raise _bad(
                    field=f"{field}.font.size",
                    expected="number >= 1",
                    received=_summary(size_pt_raw),
                    hint="Set `font.size` in points, e.g. 18.",
                )
            size_pt = int(size_pt_raw)
        if "weight" in font_d:
            weight = _check_enum(f"{field}.font.weight", font_d["weight"], _VALID_WEIGHT)
        if "italic" in font_d:
            italic_raw = font_d["italic"]
            if not isinstance(italic_raw, bool):
                raise _bad(
                    field=f"{field}.font.italic",
                    expected="boolean",
                    received=_summary(italic_raw),
                    hint="`italic` must be true or false.",
                )
            italic = italic_raw
        if "color_hex" in font_d:
            color = _check_hex(f"{field}.font.color_hex", font_d["color_hex"])

    align = _check_enum(f"{field}.align", shape.get("align", "left"), _VALID_ALIGN)
    v_align = _check_enum(f"{field}.v_align", shape.get("v_align", "top"), _VALID_V_ALIGN)
    line_spacing_raw = shape.get("line_spacing", 1.0)
    if (
        not isinstance(line_spacing_raw, int | float)
        or isinstance(line_spacing_raw, bool)
        or line_spacing_raw <= 0
    ):
        raise _bad(
            field=f"{field}.line_spacing",
            expected="positive number",
            received=_summary(line_spacing_raw),
            hint="Use values like 1.0 (default), 1.2, 1.5.",
        )
    line_spacing = float(line_spacing_raw)

    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[v_align]

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    lines = text_value.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align_map[align]
        para.line_spacing = line_spacing
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = weight == "bold"
        run.font.italic = italic
        if family:
            run.font.name = family
        if color:
            run.font.color.rgb = _rgb(color)


def _render_rect(
    slide: Any, field: str, shape: dict[str, Any], page_w: float, page_h: float
) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Pt

    x, y, w, h = _shape_bounds(field, shape, page_w, page_h)
    corner_in_raw = shape.get("corner_radius_in", 0)
    if (
        not isinstance(corner_in_raw, int | float)
        or isinstance(corner_in_raw, bool)
        or corner_in_raw < 0
    ):
        raise _bad(
            field=f"{field}.corner_radius_in",
            expected="non-negative number",
            received=_summary(corner_in_raw),
            hint="Pass a number in inches (0 = sharp corners).",
        )
    corner_in = float(corner_in_raw)
    auto_shape = MSO_SHAPE.ROUNDED_RECTANGLE if corner_in > 0 else MSO_SHAPE.RECTANGLE
    rect = slide.shapes.add_shape(
        auto_shape, _emu_inch(x), _emu_inch(y), _emu_inch(w), _emu_inch(h)
    )

    # Adjust the rounding amount to the requested radius. python-pptx uses
    # `adjustments[0]` in [0,1] = ratio of half the shorter side.
    if corner_in > 0:
        short = min(w, h) / 2 if min(w, h) > 0 else 1
        ratio = max(0.0, min(corner_in / short, 0.5))
        rect.adjustments[0] = ratio

    fill_hex = shape.get("fill_hex")
    if fill_hex is None:
        rect.fill.background()
    else:
        rect.fill.solid()
        rect.fill.fore_color.rgb = _rgb(_check_hex(f"{field}.fill_hex", fill_hex))

    border = shape.get("border")
    if border is None:
        rect.line.fill.background()
    else:
        border_d = _as_dict(border, f"{field}.border")
        rect.line.color.rgb = _rgb(
            _check_hex(f"{field}.border.color_hex", border_d.get("color_hex"))
        )
        thickness_raw = border_d.get("thickness_pt", 1)
        if (
            not isinstance(thickness_raw, int | float)
            or isinstance(thickness_raw, bool)
            or thickness_raw <= 0
        ):
            raise _bad(
                field=f"{field}.border.thickness_pt",
                expected="positive number",
                received=_summary(thickness_raw),
                hint="Set thickness in points, e.g. 1.5.",
            )
        rect.line.width = Pt(float(thickness_raw))


def _render_line(
    slide: Any, field: str, shape: dict[str, Any], page_w: float, page_h: float
) -> None:
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    from pptx.util import Pt

    x1 = _check_number(f"{field}.x1", shape.get("x1"), lo=0, hi=page_w)
    y1 = _check_number(f"{field}.y1", shape.get("y1"), lo=0, hi=page_h)
    x2 = _check_number(f"{field}.x2", shape.get("x2"), lo=0, hi=page_w)
    y2 = _check_number(f"{field}.y2", shape.get("y2"), lo=0, hi=page_h)
    color_hex = _check_hex(f"{field}.color_hex", shape.get("color_hex"))
    thickness_raw = shape.get("thickness_pt", 1)
    if (
        not isinstance(thickness_raw, int | float)
        or isinstance(thickness_raw, bool)
        or thickness_raw <= 0
    ):
        raise _bad(
            field=f"{field}.thickness_pt",
            expected="positive number",
            received=_summary(thickness_raw),
            hint="Set thickness in points, e.g. 1.5.",
        )
    dash = _check_enum(f"{field}.dash", shape.get("dash", "solid"), _VALID_DASH)

    line = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        _emu_inch(x1),
        _emu_inch(y1),
        _emu_inch(x2),
        _emu_inch(y2),
    )
    line.line.color.rgb = _rgb(color_hex)
    line.line.width = Pt(float(thickness_raw))
    if dash != "solid":
        from pptx.enum.dml import MSO_LINE_DASH_STYLE

        line.line.dash_style = (
            MSO_LINE_DASH_STYLE.DASH if dash == "dash" else MSO_LINE_DASH_STYLE.ROUND_DOT
        )


def _render_image(
    slide: Any, field: str, shape: dict[str, Any], page_w: float, page_h: float
) -> None:
    x, y, w, h = _shape_bounds(field, shape, page_w, page_h)
    data_b64 = shape.get("data_b64")
    if not isinstance(data_b64, str) or not data_b64:
        raise _bad(
            field=f"{field}.data_b64",
            expected="non-empty base64-encoded string",
            received=_summary(data_b64),
            hint="Pass `data_b64` as base64-encoded PNG or JPEG bytes (no data: prefix).",
        )
    # Strip an accidental data URL prefix the LLM may add.
    if data_b64.startswith("data:"):
        comma = data_b64.find(",")
        if comma == -1:
            raise _bad(
                field=f"{field}.data_b64",
                expected="raw base64 (no data: prefix needed)",
                received="data URL without comma separator",
                hint="Pass only the base64 chunk, e.g. 'iVBORw0KGgo...'; no leading 'data:image/png;base64,'.",
            )
        data_b64 = data_b64[comma + 1 :]
    try:
        raw = base64.b64decode(data_b64, validate=True)
    except (ValueError, binascii.Error) as exc:
        raise _bad(
            field=f"{field}.data_b64",
            expected="valid base64 string",
            received=f"{len(data_b64)} chars; decode error: {exc}",
            hint="Re-encode the image bytes as base64 (PNG or JPEG) without line breaks.",
        ) from exc
    if not raw:
        raise _bad(
            field=f"{field}.data_b64",
            expected="non-empty image bytes",
            received="0 decoded bytes",
            hint="Provide actual image bytes, not an empty string.",
        )
    if len(raw) > _MAX_IMAGE_BYTES:
        raise _bad(
            field=f"{field}.data_b64",
            expected=f"image bytes <= {_MAX_IMAGE_BYTES // 1024}KB",
            received=f"{len(raw)} bytes",
            hint="Compress the image (resize / lower JPEG quality) before embedding.",
        )

    slide.shapes.add_picture(
        io.BytesIO(raw),
        _emu_inch(x),
        _emu_inch(y),
        _emu_inch(w),
        _emu_inch(h),
    )


def _render_chart(
    slide: Any, field: str, shape: dict[str, Any], page_w: float, page_h: float
) -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    x, y, w, h = _shape_bounds(field, shape, page_w, page_h)
    chart_type = _check_enum(f"{field}.chart_type", shape.get("chart_type"), _VALID_CHART_TYPES)

    categories = shape.get("categories")
    if not isinstance(categories, list) or not categories:
        raise _bad(
            field=f"{field}.categories",
            expected="non-empty array",
            received=_summary(categories),
            hint="Pass `categories` as a list of category names (one per data point).",
        )

    series = shape.get("series")
    if not isinstance(series, list) or not series:
        raise _bad(
            field=f"{field}.series",
            expected="non-empty array of {name, values}",
            received=_summary(series),
            hint='Pass `series` as e.g. [{"name":"Q1","values":[1,2,3]}].',
        )

    chart_data = CategoryChartData()  # type: ignore[no-untyped-call]
    chart_data.categories = [str(c) for c in categories]
    series_colors: list[str | None] = []
    for s_idx, s_raw in enumerate(series):
        s_field = f"{field}.series[{s_idx}]"
        s = _as_dict(s_raw, s_field)
        if "name" not in s or not isinstance(s["name"], str) or not s["name"]:
            raise _bad(
                field=f"{s_field}.name",
                expected="non-empty string",
                received=_summary(s.get("name")),
                hint="Each series needs a `name` for the legend.",
            )
        values = s.get("values")
        if not isinstance(values, list):
            raise _bad(
                field=f"{s_field}.values",
                expected="array of numbers",
                received=_summary(values),
                hint="Pass `values` as a numeric array.",
            )
        if len(values) != len(categories):
            raise _bad(
                field=f"{s_field}.values",
                expected=f"array of length {len(categories)} (matches categories)",
                received=f"array of length {len(values)}",
                hint=(
                    "Series `values` must have the same length as `categories`. "
                    f"categories has {len(categories)} entries."
                ),
            )
        nums: list[float] = []
        for v_idx, v in enumerate(values):
            if isinstance(v, bool) or not isinstance(v, int | float):
                raise _bad(
                    field=f"{s_field}.values[{v_idx}]",
                    expected="number",
                    received=_summary(v),
                    hint="Each chart value must be a JSON number.",
                )
            nums.append(float(v))
        chart_data.add_series(s["name"], nums)  # type: ignore[no-untyped-call]
        if "color_hex" in s:
            series_colors.append(_check_hex(f"{s_field}.color_hex", s["color_hex"]))
        else:
            series_colors.append(None)

    type_map = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
    }
    chart = slide.shapes.add_chart(
        type_map[chart_type],
        _emu_inch(x),
        _emu_inch(y),
        _emu_inch(w),
        _emu_inch(h),
        chart_data,
    ).chart

    show_legend = shape.get("show_legend", True)
    if not isinstance(show_legend, bool):
        raise _bad(
            field=f"{field}.show_legend",
            expected="boolean",
            received=_summary(show_legend),
            hint="`show_legend` must be true or false.",
        )
    chart.has_legend = show_legend
    if show_legend and chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    show_data_labels = shape.get("show_data_labels", False)
    if not isinstance(show_data_labels, bool):
        raise _bad(
            field=f"{field}.show_data_labels",
            expected="boolean",
            received=_summary(show_data_labels),
            hint="`show_data_labels` must be true or false.",
        )
    if show_data_labels:
        for plot in chart.plots:
            plot.has_data_labels = True

    # Override series colors so the caller's design language wins over
    # python-pptx's default palette (chart_type=pie colors per point).
    if chart_type == "pie":
        plot = chart.plots[0]
        for pt_idx, color_hex in enumerate(series_colors):
            if color_hex is None:
                continue
            try:
                point = plot.series[0].points[pt_idx]
            except (IndexError, AttributeError):  # pragma: no cover - python-pptx safety
                continue
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _rgb(color_hex)
    else:
        for s_idx, color_hex in enumerate(series_colors):
            if color_hex is None:
                continue
            try:
                ser = chart.series[s_idx]
            except (IndexError, AttributeError):  # pragma: no cover
                continue
            if chart_type == "line":
                ser.format.line.color.rgb = _rgb(color_hex)
            else:
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = _rgb(color_hex)


def _render_one_shape(slide: Any, field: str, shape: Any, page_w: float, page_h: float) -> None:
    if not isinstance(shape, dict):
        raise _bad(
            field=field,
            expected="object",
            received=_summary(shape),
            hint=f"Each entry in `{field}` must be a JSON object literal.",
        )
    s_type = shape.get("type")
    if s_type not in _VALID_PRIMITIVES:
        raise _bad(
            field=f"{field}.type",
            expected=f"one of {list(_VALID_PRIMITIVES)}",
            received=_summary(s_type),
            hint=f"Set `{field}.type` to one of {list(_VALID_PRIMITIVES)}.",
        )
    dispatch = {
        "text": _render_text,
        "rect": _render_rect,
        "line": _render_line,
        "image": _render_image,
        "chart": _render_chart,
    }
    dispatch[s_type](slide, field, shape, page_w, page_h)


def render_pptx(
    *,
    page: dict[str, Any] | None = None,
    slides: list[Any],
) -> tuple[bytes, list[str]]:
    """Build a deck from primitive shape specs.

    Returns ``(bytes, warnings)``. The warnings list is reserved for
    non-fatal issues; the current implementation never appends to it
    because schema-level problems raise `ToolArgError` upfront — but
    keeping the channel keeps the executor / persistence contract
    stable across all office generators.
    """
    if not isinstance(slides, list) or not slides:
        raise ArtifactGenerationError("pptx requires at least one slide.")

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:  # pragma: no cover
        raise ArtifactGenerationError(f"python-pptx unavailable: {exc}") from exc

    page_d = _as_dict(page, "page") if page is not None else {}
    page_w = (
        _check_number("page.width_in", page_d.get("width_in"), lo=_MIN_DIM_IN, hi=_MAX_DIM_IN)
        if "width_in" in page_d
        else _DEFAULT_PAGE_W
    )
    page_h = (
        _check_number("page.height_in", page_d.get("height_in"), lo=_MIN_DIM_IN, hi=_MAX_DIM_IN)
        if "height_in" in page_d
        else _DEFAULT_PAGE_H
    )
    page_bg_hex: str | None = None
    if "background" in page_d:
        bg_d = _as_dict(page_d["background"], "page.background")
        if "color_hex" in bg_d:
            page_bg_hex = _check_hex("page.background.color_hex", bg_d["color_hex"])

    prs = Presentation()
    prs.slide_width = Inches(page_w)
    prs.slide_height = Inches(page_h)
    blank_layout = prs.slide_layouts[6]

    warnings: list[str] = []

    for s_idx, spec in enumerate(slides):
        s_field = f"slides[{s_idx}]"
        spec_d = _as_dict(spec, s_field)

        slide = prs.slides.add_slide(blank_layout)

        slide_bg_hex: str | None = page_bg_hex
        if "background" in spec_d:
            sb = _as_dict(spec_d["background"], f"{s_field}.background")
            if "color_hex" in sb:
                slide_bg_hex = _check_hex(f"{s_field}.background.color_hex", sb["color_hex"])
        if slide_bg_hex is not None:
            _set_slide_background(slide, slide_bg_hex)

        shapes_raw = spec_d.get("shapes")
        if not isinstance(shapes_raw, list) or not shapes_raw:
            raise _bad(
                field=f"{s_field}.shapes",
                expected="non-empty array",
                received=_summary(shapes_raw),
                hint=f"Each slide needs at least one shape · pass {s_field}.shapes as an array of primitives.",
            )
        for sh_idx, shape in enumerate(shapes_raw):
            _render_one_shape(slide, f"{s_field}.shapes[{sh_idx}]", shape, page_w, page_h)

        notes_raw = spec_d.get("notes")
        if notes_raw is not None:
            if not isinstance(notes_raw, str):
                raise _bad(
                    field=f"{s_field}.notes",
                    expected="string",
                    received=_summary(notes_raw),
                    hint="Speaker notes must be a plain string.",
                )
            if notes_raw:
                slide.notes_slide.notes_text_frame.text = notes_raw

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), warnings


def extract_slide_text(blob: bytes) -> list[dict[str, Any]]:
    """Read a .pptx blob and return ``[{title, body[]}, ...]``.

    Used by the render endpoint when the frontend pptx viewer wants a
    text outline. Pure read · no side effects. Returns ``[]`` if the
    blob is unparseable so the UI can show its 「N 张幻灯片 · 无法解析
    文本」 fallback cleanly.

    Title heuristic for primitives-mode decks: the largest top-most
    text shape on a slide is treated as the title; everything else is
    body. We sort by font size (descending) then by y-coordinate
    (ascending) and take the first non-empty as title.
    """
    try:
        from pptx import Presentation
    except ImportError:  # pragma: no cover
        return []
    try:
        prs = Presentation(io.BytesIO(blob))
    except Exception:
        return []

    out: list[dict[str, Any]] = []
    for slide in prs.slides:
        candidates: list[tuple[float, int, str]] = []
        body_lines: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            full_text = "\n".join(
                "".join(run.text for run in para.runs).strip() for para in tf.paragraphs
            ).strip()
            if not full_text:
                continue
            # Track the top y-coord and the largest font size in this text frame
            top = -float(shape.top) if shape.top is not None else 0.0
            biggest_pt = 0
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.size is not None:
                        biggest_pt = max(biggest_pt, int(run.font.size.pt))
            candidates.append((-biggest_pt, int(top), full_text))
        candidates.sort()
        title = ""
        for _, _, text in candidates:
            first = text.split("\n", 1)[0].strip()
            if first:
                title = first
                break
        for _, _, text in candidates:
            for line in text.split("\n"):
                line = line.strip()
                if line and line != title:
                    body_lines.append(line)
        out.append({"title": title, "body": body_lines})
    return out
