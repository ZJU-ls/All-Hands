"""Text extraction for non-image attachments.

Called lazily by AgentLoop's multimodal projection when a FileBlock is in
the user message: open the bytes by mime, extract up to 100KB / 50 pages
of text, return the projection-ready string. Result is cached on the
attachment row via AttachmentService.store_extracted_text so subsequent
turns don't re-extract.

Supported types:
- text/* (utf-8 decode + replace errors)
- application/json / xml / yaml (decoded as text)
- application/pdf (pypdf, first 50 pages)
- docx (python-docx, all paragraphs)
- xlsx (openpyxl, first sheet → CSV)
- pptx (python-pptx, slides → text per slide)

Anything else returns None and the caller falls back to "binary attachment,
no preview" text.
"""

from __future__ import annotations

import csv
import io
import logging
from pathlib import Path

log = logging.getLogger(__name__)

MAX_TEXT_BYTES = 100_000
MAX_PDF_PAGES = 50


def extract_text(path: Path, mime: str) -> str | None:
    """Try to extract textual content from the attachment file. Returns
    a string up to MAX_TEXT_BYTES (utf-8) or None if the type isn't
    supported / extraction fails."""
    try:
        if mime.startswith("text/") or mime in {
            "application/json",
            "application/xml",
            "application/x-yaml",
        }:
            return _extract_plain(path)
        if mime == "application/pdf":
            return _extract_pdf(path)
        if mime in {
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        }:
            return _extract_docx(path)
        if mime in {
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        }:
            return _extract_xlsx(path)
        if mime in {
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }:
            return _extract_pptx(path)
    except Exception as exc:
        log.warning("attachment.extract.failed", extra={"mime": mime, "error": str(exc)})
        return None
    return None


def _truncate(text: str) -> str:
    encoded = text.encode("utf-8")
    if len(encoded) <= MAX_TEXT_BYTES:
        return text
    return encoded[:MAX_TEXT_BYTES].decode("utf-8", errors="replace") + "\n\n[…truncated]"


def _extract_plain(path: Path) -> str:
    return _truncate(path.read_text(encoding="utf-8", errors="replace"))


def _extract_pdf(path: Path) -> str | None:
    try:
        from pypdf import PdfReader
    except ImportError:
        return None
    reader = PdfReader(str(path))
    pages = reader.pages[:MAX_PDF_PAGES]
    chunks: list[str] = []
    for i, page in enumerate(pages):
        try:
            page_text = page.extract_text() or ""
        except Exception:
            page_text = ""
        chunks.append(f"--- page {i + 1} ---\n{page_text.strip()}")
    if len(reader.pages) > MAX_PDF_PAGES:
        chunks.append(f"\n[… {len(reader.pages) - MAX_PDF_PAGES} more pages omitted]")
    return _truncate("\n\n".join(chunks))


def _extract_docx(path: Path) -> str | None:
    try:
        from docx import Document
    except ImportError:
        return None
    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text]
    return _truncate("\n".join(paragraphs))


def _extract_xlsx(path: Path) -> str | None:
    try:
        from openpyxl import load_workbook  # type: ignore[import-untyped]
    except ImportError:
        return None
    wb = load_workbook(str(path), data_only=True, read_only=True)
    out = io.StringIO()
    for sheet_name in wb.sheetnames[:5]:  # cap at 5 sheets
        out.write(f"--- sheet: {sheet_name} ---\n")
        ws = wb[sheet_name]
        writer = csv.writer(out)
        for rows_emitted, row in enumerate(ws.iter_rows(values_only=True)):
            writer.writerow(["" if v is None else str(v) for v in row])
            if rows_emitted + 1 >= 1000:
                out.write("[… more rows omitted]\n")
                break
        out.write("\n")
    wb.close()
    return _truncate(out.getvalue())


def _extract_pptx(path: Path) -> str | None:
    try:
        from pptx import Presentation
    except ImportError:
        return None
    prs = Presentation(str(path))
    chunks: list[str] = []
    for i, slide in enumerate(prs.slides[:50]):
        bits: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs)
                    if txt.strip():
                        bits.append(txt)
        chunks.append(f"--- slide {i + 1} ---\n" + "\n".join(bits))
    return _truncate("\n\n".join(chunks))
